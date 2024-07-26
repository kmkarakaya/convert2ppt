import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from matplotlib import colors as mcolors
import os



# List of valid named colors from matplotlib
valid_colors = list(mcolors.CSS4_COLORS.keys())

# Function to get valid background color from user
def get_user_background_color():

    while True:
        if os.getenv('CI'):  # or any other condition to check if running in CI
            user_specified_color = "lightblue"  # default value when in CI
        else:
            user_specified_color = input("Enter your preferred background color (e.g., 'red', 'lightblue', or HEX '#ADD8E6'): ")
        try:
            # Attempt to convert user-specified color to RGB
            rgb, rgb_tuple = convert_color(user_specified_color)
            return rgb, rgb_tuple  # Return the RGBColor object and RGB tuple
        except ValueError:
            print(f"Color '{user_specified_color}' is not recognized.")
            print("Valid colors include:")
            print(", ".join(valid_colors[:10]))  # Show first 10 valid colors for brevity
            print(f"And more. Please enter a valid color name or HEX code.")

# Convert user-specified color to RGB
def convert_color(color_name):
    try:
        rgb = mcolors.to_rgb(color_name)
    except ValueError:
        raise ValueError(f"Color '{color_name}' is not recognized.")
    
    r = int(rgb[0] * 255)
    g = int(rgb[1] * 255)
    b = int(rgb[2] * 255)
    return RGBColor(r, g, b), (r, g, b)  # Return both RGBColor and RGB components as a tuple

# Calculate brightness and choose a contrasting text color
def get_text_color(bg_rgb_tuple):
    # Extract RGB values from the tuple
    r, g, b = bg_rgb_tuple
    # Calculate luminance
    luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255  # Normalize to 0-1
    # Determine text color based on luminance
    if luminance < 0.5:  # Dark background
        return RGBColor(255, 255, 255)  # White text
    else:  # Light background
        return RGBColor(0, 0, 0)  # Black text

# Get the valid background color from user
background_color_rgb, background_color_tuple = get_user_background_color()
text_color_rgb = get_text_color(background_color_tuple)




# Utility function to add text boxes
def add_textbox(slide, text, left, top, width, height, font_size=Pt(24), bold=False, align=PP_ALIGN.LEFT, color=RGBColor(0, 0, 0)):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color  # Set the text color
    p.alignment = align
    return textbox

# Function to layout content dynamically
def layout_content(slide, body, start_top):
    current_top = start_top
    bullet_points = []
    spacing_between_elements = Inches(0.1)

    for item in body:
        if item['type'] == 'Heading':
            if bullet_points:
                add_textbox(slide,
                            "\n".join(f"• {bp}" for bp in bullet_points),
                            Inches(1), current_top, Inches(8.5), Inches(1),
                            Pt(16), color=text_color_rgb)  # Use text color
                current_top += Inches(1) + spacing_between_elements
                bullet_points = []

            add_textbox(slide,
                        item['content'],
                        Inches(1), current_top, Inches(8.5), Inches(0.5),
                        Pt(24), bold=True, color=text_color_rgb)  # Use text color
            current_top += Inches(0.5) + spacing_between_elements

        elif item['type'] == 'Bullet Point':
            bullet_points.append(item['content'])

    if bullet_points:
        add_textbox(slide,
                    "\n".join(f"• {bp}" for bp in bullet_points),
                    Inches(1), current_top, Inches(8.5), Inches(1),
                    Pt(16), color=text_color_rgb)  # Use text color
        current_top += Inches(1) + spacing_between_elements

    return current_top


# Add slides based on JSON data
def prepare_presentation(json_data, filename, data_dir):

    # Create a Presentation instance
    prs = Presentation()
    print(f"\tCreating slides for '{filename}'")
    for slide_data in json_data['slides']:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide layout

        # Apply user-specified background color
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = background_color_rgb

        title_height = Inches(1)
        add_textbox(slide, slide_data['title'], Inches(0.5), Inches(0.5), Inches(8.5), title_height, Pt(32), bold=True,
                    align=PP_ALIGN.CENTER, color=text_color_rgb)  # Use text color

        subtitle_top = Inches(0.5) + title_height
        if slide_data['subtitle']:
            subtitle_height = Inches(0.5)
            add_textbox(slide, slide_data['subtitle'], Inches(0.5), subtitle_top, Inches(8.5), subtitle_height, Pt(20),
                        align=PP_ALIGN.CENTER, color=text_color_rgb)  # Use text color
            content_top = subtitle_top + subtitle_height + Inches(0.5)
        else:
            content_top = subtitle_top + Inches(0.5)

        layout_content(slide, slide_data['body'], content_top)

        slide_number = slide_data['slide_number']
        add_textbox(slide, f"Slide {slide_number}", Inches(0.2), Inches(7.5), Inches(1), Inches(0.3), Pt(10), color=text_color_rgb)  # Use text color

    # Save presentation
    filename = filename.replace('.json', '.pptx')
    full_path = os.path.join(data_dir, filename)
    print(f"\tSaving presentation to '{full_path}'")
    prs.save(full_path)


def create_presentations(data_dir):
    # Read a JSON object from the local directory every .json file
    print("Starting building PPT file(s).")
    for file_name in os.listdir(data_dir):
        if file_name.endswith('.json'):
            with open(os.path.join(data_dir, file_name), 'r', encoding='utf-8') as file:
                print(f"\tLoading Json file: {file_name}")
                json_data = json.load(file)
                prepare_presentation(json_data, file_name,data_dir)
                print(f"\tThe PPT file for {file_name} created successfully: ")
    print("All presentations files created successfully.")
    
